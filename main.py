from fastapi import FastAPI, UploadFile, Form,File,Request
from fastapi.middleware.cors import CORSMiddleware
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core import Document, GPTListIndex, VectorStoreIndex
from llama_index.core import StorageContext, load_index_from_storage
from fastapi import HTTPException
from fastapi import Query
import mimetypes
import PyPDF2
from llama_index.core.retrievers import VectorIndexRetriever 
from llama_index.core.query_engine import RetrieverQueryEngine
from llama_index.core import PromptTemplate
from llama_index.core.indices.postprocessor import SimilarityPostprocessor
from llama_index.core.response_synthesizers import get_response_synthesizer
from typing import List
from pptx import Presentation
from io import StringIO
import os 
from io import BytesIO
import shutil
import uuid
import random
from llama_index.llms.azure_openai import AzureOpenAI
from llama_index.core import PromptTemplate

from dotenv import load_dotenv
load_dotenv()

os.environ["AZURE_API_KEY"] = os.getenv("AZURE_API_KEY") # azure
os.environ["AZURE_ENDPOINT"] = (
    os.getenv("AZURE_ENDPOINT")  # azure
)
os.environ["AZURE_API_VERSION"] = os.getenv("AZURE_API_VERSION")  # may need updating
os.environ["OPENAI_API_TYPE"] = "azure"

TEMP_DIR_BASE = "powerpoint"

def generate_temp_directory() -> str:
    """Generates a unique directory path and creates the directory."""
    unique_dir = os.path.join(TEMP_DIR_BASE, str(uuid.uuid4()))
    os.makedirs(unique_dir, exist_ok=True)
    return unique_dir

def cleanup_temp_directory(directory: str):
    """Cleans up the temporary directory and its contents."""
    shutil.rmtree(directory)

app = FastAPI()

def load_collection():
    #vector_store = FaissVectorStore.from_persist_dir(persist_dir="./storage")
    
    storage_context = StorageContext.from_defaults(
         persist_dir="./temp_storage"
    )

    index = load_index_from_storage(storage_context=storage_context)

    return index

# Allow CORS for frontend communication
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Global index initialized
index = None
# Global dict_name to remember which document to delete for the delete function
dict_name = {}
node_list = []

uploadCount = 0

@app.post("/upload")
async def upload_document(files: List[UploadFile]):
    global index
    global dict_name
    global uploadCount
    global node_list
    try:
        documents = []
        i=0
        # Check file type
        for file in files:
            mime_type, _ = mimetypes.guess_type(file.filename)
            if mime_type == "application/pdf":
                # Extract text from PDF
                pdf_reader = PyPDF2.PdfReader(file.file)
                content = "\n".join(page.extract_text() for page in pdf_reader.pages if page.extract_text())
            elif mime_type == "text/markdown":
                # Handle markdown files
                content = (await file.read()).decode("utf-8")
            elif mime_type == "text/plain":
                # Handle plain text files
                content = (await file.read()).decode("utf-8")
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":

                content = await file.read()
                
                prs = Presentation(BytesIO(content))
                text_content = "\n".join([slide.notes_slide.notes_text_frame.text.strip() 
                                          for slide in prs.slides if slide.notes_slide and slide.notes_slide.notes_text_frame])
                
                if not text_content:
                    text_content = "\n".join([shape.text.strip() for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                
                content = text_content.encode("utf-8")  # Convert to UTF-8 encoded text
                
            else:
                raise HTTPException(status_code=400, detail="Unsupported file type.")
            
            # Create and index the document
            document = Document(text=content,id_=f'doc_id:{i}')
            dict_name[file.filename] = f'doc_id:{i}'
            i+=1
            documents.append(document)

            splitter = SentenceSplitter(chunk_size=2048,chunk_overlap=200) #splitting the data
            nodes = splitter.get_nodes_from_documents(documents)
            node_list = node_list + list(nodes)
            #print(len(nodes))
            #print(len(node_list))
            #print("No of nodes: ", len(nodes))
            #print("Ref doc id", nodes[0].ref_doc_id)
        if uploadCount ==0:
            index = VectorStoreIndex(nodes) #GPTListIndex(nodes)
            index.storage_context.persist("./temp_storage")
        else:
            index.insert_nodes(nodes)
        uploadCount+=1
        return {"message": "Document(s) uploaded and indexed successfully."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to process the file: {str(e)}")

@app.post("/query")
async def query_index(query: str = File(...)):
    global index
    if not index:
        return {"error": "No document has been uploaded yet."}
    
    #query_engine = index.as_query_engine()
    #response = query_engine.query(query)
    retriever = VectorIndexRetriever(index=index, similarity_top_k=4) #modify the index to use 4 source nodes 
    postprocessor = SimilarityPostprocessor(similarity_cutoff=0.7)

    system_prompt = PromptTemplate(
    "You are a helpful assistant. Use the context below to answer the query concisely and accurately. Otherwise, simply say 'I don't know'.\n\n{context_str}\n\nQuestion: {query_str}\nAnswer:"
    )

    response_synthesizer = get_response_synthesizer(simple_template=system_prompt)

    query_engine = RetrieverQueryEngine(retriever=retriever,
                                        node_postprocessors=[postprocessor],
                                        response_synthesizer=response_synthesizer
                                        )
    
    response = query_engine.query(query)

    return {"response": response.response}

@app.delete("/delete")
async def delete_document(doc_name:str = Query(...)):#request: Request):
    global index, dict_name
    #print(dict_name)
    #print(index.docstore.get_all_ref_doc_info())
    try:
        '''body = await request.json()  # Parse the JSON body
        doc_name = body.get("doc_name")

        if not doc_name:
            raise HTTPException(status_code=400, detail="Document name is required.")
        if doc_name not in dict_name:
            raise HTTPException(status_code=404, detail="Document not found.")'''

        if doc_name not in dict_name:
            raise HTTPException(status_code=404, detail="Document not found.")
        
        doc_id = dict_name[doc_name]

        # Delete the document by its ID
        index.delete_ref_doc(doc_id, delete_from_docstore=True)
        index.storage_context.persist('./temp_storage')  # Persist the changes

        # Remove the document from the dictionary
        del dict_name[doc_name]

        index = load_collection()
        #print(index.docstore.get_all_ref_doc_info())
        #print(dict_name)

        return {"message": f"Document '{doc_name}' deleted successfully."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to delete document: {str(e)}")
    

@app.get("/list_uploaded_files")
async def list_uploaded_files():
    global dict_name
    return {"files": list(dict_name.keys())}


def create_nodes(n=3):
    global node_list
    nodes = node_list

    #print("The no. of nodes are: ", len(nodes))

    selected_nodes = random.sample(nodes, n)

    #print('Error in create_nodes?')
    
        #raise ValueError("Not enough source nodes to select from.")
    
    separator = "----------"
    context = f"\n {separator} \n".join(node.get_text() for node in selected_nodes)

    return context

def initialize_llm():

    #print('error in llm?')
    llm = AzureOpenAI(
    model="gpt-4o-mini",
    deployment_name="gpt-4o-mini",
    api_key=os.environ["AZURE_API_KEY"],
    azure_endpoint=os.environ["AZURE_ENDPOINT"],
    api_version=os.environ["AZURE_API_VERSION"],
    )


    return llm

def create_prompt(context,query):
    #print('error in prompt?')
    system_prompt_template = """
    You are a highly intelligent question generator. Your task is to generate unique and insightful questions, that would be asked by a human user, based on the provided data context.
    Produce at least 1 question for each passage that is separated by the '----------'.

    Context:
    {context}

    Instructions:
    1.Generate questions that are highly relevant and unique to the above context.
    2.Do not number the questions.

    {query}
    """
    qa_template = PromptTemplate(system_prompt_template)

    prompt = qa_template.format(context=context,query=query)

    return prompt

@app.get("/suggestive_questions")
async def generate_q():
    lst_filter = []
    context = create_nodes()
    query = "Generate 3 unique questions that are short and concise."

    prompt = create_prompt(context, query)

    llm = initialize_llm()

   # print('error in creating response?')
    response = llm.complete(prompt)

    lst = str(response).strip().split('\n')

    for i in lst:
        if i != "":
            lst_filter.append(i)

    return lst_filter